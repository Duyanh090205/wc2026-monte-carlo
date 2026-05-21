"""Generate Stream 3 preliminary report deck for Sam.

Plain-language content (per /explain skill principles): every technical term
glossed inline on first use; numbers carry context; uses analogies where they
help. Native python-pptx shapes for roadmap + architecture (no images).

Output: mc_simu/preliminary_report.pptx (5 slides, ~70-100KB).
"""
from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "src"))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402

DATA_DIR = PROJECT_ROOT / "data" / "mc_simu" / "phase3_baselines"
OUT_PATH = PROJECT_ROOT / "mc_simu" / "preliminary_report.pptx"

# ── Color palette ─────────────────────────────────────────────────────────────
COL_NAVY = RGBColor(0x1F, 0x2A, 0x44)         # title text
COL_INK = RGBColor(0x33, 0x33, 0x33)          # body text
COL_MUTED = RGBColor(0x88, 0x88, 0x88)
COL_GREEN = RGBColor(0x2E, 0x8B, 0x57)        # done / BUY
COL_RED = RGBColor(0xC0, 0x39, 0x2B)          # FADE / warning
COL_AMBER = RGBColor(0xE6, 0xA8, 0x17)        # in progress
COL_GRAY = RGBColor(0xCC, 0xCC, 0xCC)         # skipped / pending
COL_BG = RGBColor(0xF6, 0xF7, 0xF9)
COL_ACCENT = RGBColor(0x2A, 0x5C, 0xAA)       # eloratings (primary)
COL_ACCENT_LIGHT = RGBColor(0x6C, 0x8E, 0xBE) # self-elo (secondary)
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_text(tf, text, *, size=14, bold=False, color=COL_INK, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def _add_textbox(slide, x, y, w, h, text, *, size=14, bold=False,
                 color=COL_INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    _set_text(tf, text, size=size, bold=bold, color=color, align=align)
    return box


def _add_bullets(slide, x, y, w, h, lines, *, size=14, color=COL_INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Each line is (bullet_text, optional_bold)
        if isinstance(line, tuple):
            text, is_bold = line
        else:
            text, is_bold = line, False
        r = p.add_run()
        r.text = "•  " + text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = is_bold
        p.space_after = Pt(6)
    return box


def _box(slide, x, y, w, h, fill_color, *, line_color=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def _arrow(slide, x1, y1, x2, y2, color=COL_MUTED):
    """Down/right arrow connector."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    # Add arrowhead at end
    from pptx.oxml.ns import qn
    line = conn.line._get_or_add_ln()
    tail = line.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(line, qn("a:tailEnd"))
    tail.set("type", "triangle")
    return conn


# ── Slide 1: Title ────────────────────────────────────────────────────────────


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Decorative top stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0),
                                     prs.slide_width, Inches(0.5))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = COL_NAVY
    stripe.line.fill.background()

    _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2),
                 "Stream 3: Tournament Simulator",
                 size=44, bold=True, color=COL_NAVY)
    _add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.8),
                 "Preliminary status report — what we built, what we found",
                 size=22, color=COL_INK)
    _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
                 "2026-05-20  ·  Phase 0–3 complete  ·  Phase 4 (tuning) in progress",
                 size=14, color=COL_MUTED)
    _add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
                 "Audience: Sam + team review",
                 size=14, color=COL_MUTED)


# ── Slide 2: Why + Roadmap ────────────────────────────────────────────────────


def slide_why_and_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 "Why we need a third signal — and where we are",
                 size=28, bold=True, color=COL_NAVY)

    # Why section (top half)
    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
                 "Problem", size=16, bold=True, color=COL_ACCENT)
    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.2),
                 "FairLine already has 2 ways to estimate \"who will win the tournament\": "
                 "sportsbook odds (Stream 1) and prediction market prices (Stream 2). "
                 "Both reflect what the market thinks. When the market is collectively wrong "
                 "in the same direction, neither stream catches it.",
                 size=14, color=COL_INK)

    _add_textbox(slide, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.4),
                 "Stream 3 = independent third source", size=16, bold=True, color=COL_ACCENT)
    _add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.6),
                 "Predicts champions from team strength (Elo ratings) + tournament bracket "
                 "structure — NOT from market prices. Disagreement with Streams 1/2 is the signal.",
                 size=14, color=COL_INK)

    # Roadmap (bottom half)
    _add_textbox(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.4),
                 "Roadmap", size=16, bold=True, color=COL_ACCENT)

    phases = [
        ("Phase 0", "Data audit", "1d", COL_GREEN, "DONE"),
        ("Phase 1", "Build Elo +\nsingle-game", "1d", COL_GREEN, "DONE"),
        ("Phase 2", "Roster Elo", "—", COL_GRAY, "SKIPPED"),
        ("Phase 3", "Tournament\nsimulator", "1d", COL_GREEN, "DONE"),
        ("Phase 4", "Tune +\ncompare", "now", COL_AMBER, "IN PROGRESS"),
        ("Phase 5", "Present to\nteam", "—", COL_GRAY, "PENDING"),
    ]
    n = len(phases)
    margin = Inches(0.5)
    gap = Inches(0.18)
    avail = prs.slide_width - 2 * margin - gap * (n - 1)
    box_w = avail // n
    box_h = Inches(1.4)
    y0 = Inches(4.5)

    for i, (name, desc, days, color, status) in enumerate(phases):
        x = margin + i * (box_w + gap)
        # Main box
        _box(slide, x, y0, box_w, box_h, color)
        # Phase name (bold, top)
        tb = slide.shapes.add_textbox(x, y0 + Inches(0.1),
                                       box_w, Inches(0.35))
        _set_text(tb.text_frame, name, size=14, bold=True,
                  color=COL_WHITE, align=PP_ALIGN.CENTER)
        # Description (middle)
        tb2 = slide.shapes.add_textbox(x, y0 + Inches(0.45),
                                        box_w, Inches(0.6))
        tb2.text_frame.word_wrap = True
        _set_text(tb2.text_frame, desc, size=11,
                  color=COL_WHITE, align=PP_ALIGN.CENTER)
        # Status (bottom)
        tb3 = slide.shapes.add_textbox(x, y0 + box_h - Inches(0.35),
                                        box_w, Inches(0.3))
        _set_text(tb3.text_frame, f"{status}  ·  {days}", size=10,
                  bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
        # Arrow to next phase
        if i < n - 1:
            ax1 = x + box_w
            ax2 = x + box_w + gap
            ay = y0 + box_h // 2
            _arrow(slide, ax1, ay, ax2, ay, color=COL_MUTED)

    _add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                 "Target ship: June 11 (before WC 2026 kickoff)",
                 size=12, color=COL_MUTED, align=PP_ALIGN.CENTER)


# ── Slide 3: Methodology + architecture ───────────────────────────────────────


def slide_methodology(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 "How the simulator works — Sam's plug-in framework",
                 size=28, bold=True, color=COL_NAVY)

    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.7),
                 "Three layers. Top two are FIXED across all tournaments. "
                 "Bottom layer is INTERCHANGEABLE — we plug in different prediction "
                 "models and compare them on the same simulator.",
                 size=13, color=COL_INK)

    # Architecture diagram — left side (3 stacked boxes + 3 swappable bottom)
    arch_x = Inches(0.5)
    arch_y = Inches(2.0)
    arch_w = Inches(6.5)

    layer_h = Inches(0.95)
    gap_v = Inches(0.35)

    # Layer 1: Harness
    _box(slide, arch_x, arch_y, arch_w, layer_h, COL_ACCENT)
    _add_textbox(slide, arch_x, arch_y + Inches(0.08), arch_w, Inches(0.35),
                 "Tournament Engine (fixed)",
                 size=14, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, arch_x, arch_y + Inches(0.45), arch_w, Inches(0.45),
                 "Runs group stage matches, applies FIFA tiebreakers,\nadvances winners through knockouts to the final",
                 size=10, color=COL_WHITE, align=PP_ALIGN.CENTER)

    # Layer 2: Bracket adapter
    y2 = arch_y + layer_h + gap_v
    _box(slide, arch_x, y2, arch_w, layer_h, COL_ACCENT_LIGHT)
    _add_textbox(slide, arch_x, y2 + Inches(0.08), arch_w, Inches(0.35),
                 "Bracket Adapter (per tournament format)",
                 size=14, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, arch_x, y2 + Inches(0.45), arch_w, Inches(0.45),
                 "Defines groups + knockout seeding rules.\nDifferent for WC 2026 vs WC 2002–2022 vs Euro",
                 size=10, color=COL_WHITE, align=PP_ALIGN.CENTER)

    # Layer 3: Match predictor (swappable)
    y3 = y2 + layer_h + gap_v
    _box(slide, arch_x, y3, arch_w, Inches(1.6),
         RGBColor(0xEC, 0xEF, 0xF5), line_color=COL_ACCENT)
    _add_textbox(slide, arch_x, y3 + Inches(0.08), arch_w, Inches(0.35),
                 "Match Predictor — SWAPPABLE",
                 size=14, bold=True, color=COL_ACCENT, align=PP_ALIGN.CENTER)

    # 3 inner boxes side by side
    inner_y = y3 + Inches(0.5)
    inner_h = Inches(1.0)
    inner_gap = Inches(0.15)
    inner_margin = Inches(0.2)
    avail_inner = arch_w - 2 * inner_margin - inner_gap * 2
    inner_w = avail_inner // 3
    inner_x0 = arch_x + inner_margin

    models = [
        ("Self-Elo",  "Built from\n49k matches\n(1872–2026)", COL_ACCENT),
        ("Eloratings.net", "Scraped 244\nteams' ratings", COL_ACCENT_LIGHT),
        ("Baseline\n40/40/20", "Pure ignorance\n(sanity floor)", COL_MUTED),
    ]
    for i, (name, sub, col) in enumerate(models):
        ix = inner_x0 + i * (inner_w + inner_gap)
        _box(slide, ix, inner_y, inner_w, inner_h, col)
        _add_textbox(slide, ix, inner_y + Inches(0.1), inner_w, Inches(0.35),
                     name, size=12, bold=True, color=COL_WHITE,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, ix, inner_y + Inches(0.45), inner_w, Inches(0.55),
                     sub, size=9, color=COL_WHITE, align=PP_ALIGN.CENTER)

    # Arrows between layers
    cx = arch_x + arch_w // 2
    _arrow(slide, cx, arch_y + layer_h, cx, y2, color=COL_MUTED)
    _arrow(slide, cx, y2 + layer_h, cx, y3, color=COL_MUTED)

    # Right side: key rules
    rx = Inches(7.4)
    ry = Inches(2.0)
    _add_textbox(slide, rx, ry, Inches(5.5), Inches(0.4),
                 "Key rules", size=16, bold=True, color=COL_ACCENT)
    _add_bullets(slide, rx, ry + Inches(0.5), Inches(5.5), Inches(4.5),
                 [
                  ("Tournament engine never changes — same code runs WC 2026 + 11 past tournaments.", True),
                  "Same predictor interface for all 3 models: input (team A, team B, context), output (P_A, P_draw, P_B).",
                  ("Never tune the model to match market prices.", True),
                  "  If we did, Stream 3 would just echo Streams 1+2 → no independent signal.",
                  ("Tune only on actual scoreboard outcomes from past matches.", True),
                  "  Train on what really happened, not what the market thought.",
                  "Baseline 40/40/20 = floor: any real model must beat it. If not, model is broken.",
                 ], size=11, color=COL_INK)


# ── Slide 4: Historical replay table (plain-language version) ─────────────────


def slide_historical_replay(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.55),
                 "Did the model work? Replayed 11 past tournaments",
                 size=28, bold=True, color=COL_NAVY)

    _add_textbox(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5),
                 "For each past tournament, we asked each model: \"How likely is the team that ACTUALLY won?\" "
                 "Higher % = the model gave that champion better odds before the tournament.",
                 size=13, color=COL_INK)

    # ── Main comparison table: per-edition champion probability across 3 models ──
    # Chronological order. Each row: edition | actual winner | self | eloratings | baseline
    # Best model per row gets green highlight + bold.
    rows = [
        # (edition, winner, self_pct, elor_pct, base_pct, best_col_index)
        ("WC 2002",    "Brazil",     3.6,  3.2,  3.5),
        ("Euro 2004",  "Greece",     0.9,  1.3,  6.1),
        ("WC 2006",    "Italy",      5.7,  7.0,  3.2),
        ("Euro 2008",  "Spain",     12.4, 12.8,  6.3),
        ("WC 2010",    "Spain",     19.1, 19.6,  3.2),
        ("Euro 2012",  "Spain",     32.5, 31.3,  5.9),
        ("WC 2014",    "Germany",   13.9, 15.4,  3.2),
        ("WC 2018",    "France",     6.3,  6.9,  3.5),
        ("Euro 2020",  "Italy",      9.4,  9.3,  4.2),
        ("WC 2022",    "Argentina", 18.1, 18.9,  3.5),
        ("Euro 2024",  "Spain",     11.9, 12.4,  4.0),
    ]

    header = ["Tournament", "Actual winner",
              "Self-built Elo", "Eloratings.net", "Baseline 40/40/20"]
    table_x = Inches(0.5)
    table_y = Inches(1.6)
    table_w = Inches(12.3)
    table_h = Inches(3.55)

    n_rows = len(rows) + 1
    n_cols = 5
    table_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y,
                                          table_w, table_h)
    table = table_shape.table

    # Column widths (proportional)
    col_widths = [Inches(1.6), Inches(2.0), Inches(2.9), Inches(2.9), Inches(2.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row styling
    for c, txt in enumerate(header):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = COL_ACCENT
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = txt
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = COL_WHITE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)

    # Data rows
    for ri, (edition, winner, s, e, b) in enumerate(rows, start=1):
        probs = [s, e, b]
        best_idx = probs.index(max(probs))  # 0=self, 1=elor, 2=baseline
        col_vals = [edition, winner, f"{s:.1f}%", f"{e:.1f}%", f"{b:.1f}%"]
        for c, val in enumerate(col_vals):
            cell = table.cell(ri, c)
            # Background: best-model cell gets light green; others white
            if c >= 2 and (c - 2) == best_idx:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xDC, 0xEF, 0xDB)  # light green
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = COL_WHITE
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = val
            r.font.size = Pt(11)
            # Bold the actual-winner column + bold best-model cell
            if c == 1 or (c >= 2 and (c - 2) == best_idx):
                r.font.bold = True
            if c == 1:
                r.font.color.rgb = COL_NAVY
            elif c >= 2 and (c - 2) == best_idx:
                r.font.color.rgb = COL_GREEN
            else:
                r.font.color.rgb = COL_INK
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)

    # Legend
    _add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.35),
                 "Green cell = model that gave the actual champion the highest probability that tournament.",
                 size=10, color=COL_MUTED)

    # Plain-language summary
    _add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
                 "What the numbers say",
                 size=15, bold=True, color=COL_ACCENT)
    _add_bullets(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(1.3),
                 [
                  ("Eloratings.net gave the actual champion the best odds in 7 of 11 tournaments. "
                   "Self-built Elo won 3 times. Baseline won 1 — the Greece 2004 upset, by accident "
                   "(baseline gives every team ~3–6%, so it sometimes \"lucks into\" an underdog winner).", False),
                  ("Both real models put the eventual champion in their top 3 about 55% of the time — "
                   "meaningfully better than random.", False),
                  ("Eloratings and self-Elo are very close. The lead may flip after Phase 4 tuning.", True),
                 ], size=11, color=COL_INK)


# ── Slide 5: WC 2026 vs Polymarket ────────────────────────────────────────────


def slide_wc2026_vs_market(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 "WC 2026 outlook — model vs Polymarket",
                 size=28, bold=True, color=COL_NAVY)

    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.7),
                 "Ran 100,000 simulated World Cups using the eloratings model, then compared "
                 "our champion-probability for each team against Polymarket prices. Differences "
                 "of ≥3 percentage points (\"pp\") flagged as potential mispricings.",
                 size=13, color=COL_INK)

    # ── Side-by-side top 10 ──
    _add_textbox(slide, Inches(0.5), Inches(1.95), Inches(6.2), Inches(0.4),
                 "Top 10 contenders — model vs market",
                 size=14, bold=True, color=COL_ACCENT)
    # "—" means our PM data pull didn't include this market (data gap), NOT
    # that the market is untraded. Polymarket has markets for all WC2026 teams.
    rows_a = [
        ["Team", "Polymarket", "Model (eloratings)"],
        ["Spain",        "15.6%", "20.8%"],
        ["Argentina",     "8.9%", "14.0%"],
        ["France",       "16.7%", "12.5%"],
        ["England",      "11.2%", "6.8%"],
        ["Brazil",        "8.6%", "4.5%"],
        ["Portugal",          "—", "4.5%"],
        ["Colombia",          "—", "3.9%"],
        ["Netherlands",       "—", "4.0%"],
        ["Ecuador",           "—", "2.9%"],
        ["Germany",           "—", "2.8%"],
    ]
    _make_table(slide, Inches(0.5), Inches(2.4), Inches(6.2), Inches(3.3),
                rows_a, header_fill=COL_ACCENT, font_size_data=10,
                col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
    _add_textbox(slide, Inches(0.5), Inches(5.75), Inches(6.2), Inches(0.35),
                 "\"—\" = market exists on Polymarket but not in our data pull (top-5 liquidity only).",
                 size=9, color=COL_MUTED)

    # ── Edges ──
    _add_textbox(slide, Inches(7.0), Inches(1.95), Inches(5.8), Inches(0.4),
                 "5 actionable edges (≥3pp difference)",
                 size=14, bold=True, color=COL_ACCENT)
    edges = [
        ["Direction", "Team", "PM", "Model", "Edge"],
        ["BUY",   "Spain",     "15.6%", "20.8%", "+5.2pp"],
        ["BUY",   "Argentina",  "8.9%", "14.0%", "+5.1pp"],
        ["FADE",  "England",   "11.2%",  "6.8%", "−4.4pp"],
        ["FADE",  "France",    "16.7%", "12.5%", "−4.2pp"],
        ["FADE",  "Brazil",     "8.6%",  "4.5%", "−4.1pp"],
    ]
    _make_table(slide, Inches(7.0), Inches(2.4), Inches(5.8), Inches(2.5),
                edges, header_fill=COL_ACCENT, font_size_data=11,
                col_aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT,
                            PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
                direction_col=0)

    # ── Caveat box (red-tinted, bottom right) ──
    cv_x = Inches(7.0)
    cv_y = Inches(5.05)
    cv_w = Inches(5.8)
    cv_h = Inches(1.85)
    cv_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cv_x, cv_y, cv_w, cv_h)
    cv_box.fill.solid(); cv_box.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEA)
    cv_box.line.color.rgb = COL_RED; cv_box.line.width = Pt(1.5)
    _add_textbox(slide, cv_x + Inches(0.15), cv_y + Inches(0.1),
                 cv_w - Inches(0.3), Inches(0.4),
                 "⚠  DO NOT TRADE on these numbers",
                 size=13, bold=True, color=COL_RED)
    _add_bullets(slide, cv_x + Inches(0.15), cv_y + Inches(0.5),
                 cv_w - Inches(0.3), cv_h - Inches(0.6),
                 [
                  "Model is UNTUNED — Phase 4 tuning will shift each % by ±2–3pp",
                  "Our PM data pull only covered the top-5 liquidity markets — Portugal, Colombia, Netherlands, etc. have markets we didn't capture",
                  "No transaction fees factored in (Polymarket fee eats ~1–2pp)",
                 ], size=10, color=COL_INK)

    # ── Bottom-left next steps ──
    _add_textbox(slide, Inches(0.5), Inches(6.1), Inches(6.2), Inches(0.4),
                 "What comes next (Phase 4)", size=14, bold=True, color=COL_ACCENT)
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(6.2), Inches(0.6),
                 "Tune 3 model parameters using leave-one-tournament-out\n"
                 "cross-validation → final probabilities + trade-ready output.",
                 size=11, color=COL_INK)


# ── Table helper ──────────────────────────────────────────────────────────────


def _make_table(slide, x, y, w, h, rows, *,
                header_fill=None, font_size_data=11, font_size_header=11,
                col_aligns=None, highlight_row=None, direction_col=None):
    """Create a styled table. rows[0] is header.

    direction_col: index of "Direction" column (BUY/FADE) — colored.
    highlight_row: 1-indexed row to highlight (bold + slight bg).
    """
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table

    for c, header_text in enumerate(rows[0]):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill or COL_NAVY
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = (col_aligns[c] if col_aligns else PP_ALIGN.LEFT)
        r = p.add_run(); r.text = header_text
        r.font.size = Pt(font_size_header); r.font.bold = True
        r.font.color.rgb = COL_WHITE
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)

    for ri, row_data in enumerate(rows[1:], start=1):
        is_highlight = (highlight_row == ri)
        for c, text in enumerate(row_data):
            cell = table.cell(ri, c)
            if is_highlight:
                cell.fill.solid(); cell.fill.fore_color.rgb = COL_BG
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = COL_WHITE
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = (col_aligns[c] if col_aligns else PP_ALIGN.LEFT)
            r = p.add_run(); r.text = str(text)
            r.font.size = Pt(font_size_data)
            r.font.bold = is_highlight or (direction_col is not None and c == direction_col)
            # Color BUY green, FADE red
            if direction_col is not None and c == direction_col:
                if str(text).upper() == "BUY":
                    r.font.color.rgb = COL_GREEN
                elif str(text).upper() == "FADE":
                    r.font.color.rgb = COL_RED
            else:
                r.font.color.rgb = COL_INK
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)


# ── Main ──────────────────────────────────────────────────────────────────────


def build() -> Path:
    prs = Presentation()
    # 16:9 widescreen 13.33" × 7.5"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_why_and_roadmap(prs)
    slide_methodology(prs)
    slide_historical_replay(prs)
    slide_wc2026_vs_market(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    p = build()
    size_kb = p.stat().st_size / 1024
    print(f"Wrote {p}  ·  {size_kb:.1f} KB  ·  5 slides")
